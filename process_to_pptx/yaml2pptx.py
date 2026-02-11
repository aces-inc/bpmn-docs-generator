"""YAML 業務プロセス定義から編集可能な PPTX を生成する。"""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

from .yaml_loader import (
    ProcessLayout,
    load_process_yaml,
    compute_layout,
    EMU_PER_PT,
)


def _add_arrow_to_connector(connector) -> None:
    """コネクタの終端（接続先側）に矢印を付ける。headEnd が end_connect 側。"""
    line_elem = connector.line._get_or_add_ln()
    line_elem.append(
        parse_xml(
            '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
        )
    )


def _set_connector_dotted(connector) -> None:
    """コネクタを点線にする。"""
    ln = connector.line._get_or_add_ln()
    dash = parse_xml(
        '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dot"/>'
    )
    ln.append(dash)


def _set_connector_ends(connector, tail_oval: bool = False, head_arrow: bool = False) -> None:
    """コネクタの端点を設定。tail=始点（begin_connect）、head=終点（end_connect）。"""
    ln = connector.line._get_or_add_ln()
    if tail_oval:
        ln.append(
            parse_xml(
                '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="oval" w="med" len="med"/>'
            )
        )
    if head_arrow:
        ln.append(
            parse_xml(
                '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
            )
        )


# DoD: アクター名の四角 — 点線から 2pt 離して長方形、等間隔
ACTOR_BOX_GAP_PT = 2
ACTOR_BOX_GAP_EMU = ACTOR_BOX_GAP_PT * EMU_PER_PT


def _draw_actor_labels(slide, layout: ProcessLayout) -> None:
    """スライド左側にアクター名を点線から2pt離した長方形内に描画。DoD: 角のある四角・塗りつぶしなし・枠線黒・フォント黒・影なし。"""
    for i, name in enumerate(layout.actors):
        lane_top = layout.content_top_offset + i * layout.lane_height
        # 点線の上下 2pt ずつ離して四角があり等間隔（DoD）
        top = lane_top + ACTOR_BOX_GAP_EMU
        box_height = layout.lane_height - 2 * ACTOR_BOX_GAP_EMU
        left = layout.left_margin
        width = layout.left_label_width  # アクター列幅に準拠
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,  # 角のある長方形（角丸ではない）
            Emu(left), Emu(top), Emu(width), Emu(box_height),
        )
        rect.fill.background()  # 塗りつぶしなし
        rect.line.color.rgb = RGBColor(0, 0, 0)  # 枠線黒
        rect.shadow.inherit = False  # 影なし
        tf = rect.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # テキストは上下中央
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)  # フォント黒
        p.alignment = PP_ALIGN.CENTER  # 横方向も中央


def _draw_lane_separators(slide, layout: ProcessLayout) -> None:
    """レーン間をグレーの点線で区切る。点線はレーンの左端まで届く。影なし（DoD）。"""
    gray = RGBColor(0x80, 0x80, 0x80)
    x1 = layout.left_margin  # スライド左端から 10pt 余白の内側
    x2 = int(layout.slide_width - layout.right_margin)
    for i in range(1, len(layout.actors)):
        y = layout.content_top_offset + i * layout.lane_height
        line = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, x1, y, x2, y
        )
        line.line.color.rgb = gray
        line.line.width = Pt(0.5)
        line.shadow.inherit = False  # 影なし
        # 点線: dashType を設定（a:prstDash）
        ln = line.line._get_or_add_ln()
        dash = parse_xml(
            '<a:prstDash xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="dot"/>'
        )
        ln.append(dash)


# 四角形の接続点: 0=上, 1=左, 2=下, 3=右（各辺の中央）
CONNECTION_SITE_TOP = 0
CONNECTION_SITE_LEFT = 1
CONNECTION_SITE_BOTTOM = 2
CONNECTION_SITE_RIGHT = 3


def _connection_site_from(from_node, to_node) -> int:
    """始点（from）側の接続辺: 次タスクが同レーンまたは下は右、上は上（DoD）。"""
    from_actor = from_node.actor_index
    to_actor = to_node.actor_index
    # 次のタスクが上のレーン → 上から出す
    if to_actor < from_actor:
        return CONNECTION_SITE_TOP
    # 同レーンまたは下のレーン → 右から出す
    return CONNECTION_SITE_RIGHT


def _connection_site_to(from_node, to_node) -> int:
    """終点（to）側の接続辺: 前が同レーンは左、上レーンは上、下レーンは下（DoD）。"""
    from_actor = from_node.actor_index
    to_actor = to_node.actor_index
    if from_actor < to_actor:
        return CONNECTION_SITE_TOP  # 前が上レーン → 上で受ける
    if from_actor > to_actor:
        return CONNECTION_SITE_BOTTOM  # 前が下レーン → 下で受ける
    return CONNECTION_SITE_LEFT  # 同レーン → 左で受ける


def _draw_node_shape(slide, node, left: int, top: int, width: int, height: int):
    """タスク・分岐・スタート・ゴール・成果物・サービスの図形を 1 つ描画。"""
    if node.type == "gateway":
        shape_type = MSO_SHAPE.DIAMOND
    elif node.type in ("start", "end"):
        shape_type = MSO_SHAPE.OVAL  # 正円は幅＝高さの楕円で描画
    elif node.type == "artifact":
        shape_type = MSO_SHAPE.FLOWCHART_DATA  # 成果物＝フローチャートのデータ図形（DoD）
    elif node.type == "service":
        shape_type = MSO_SHAPE.FLOWCHART_MAGNETIC_DISK  # システム接続のサービス（DoD）
    else:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

    # スタート・ゴールは正円のため、セル内で幅＝高さにし中央に配置
    if node.type in ("start", "end"):
        side = min(width, height)
        left = left + (width - side) // 2
        top = top + (height - side) // 2
        width = height = side
    # サービスはタスクと同サイズの磁気ディスク
    if node.type == "service":
        side = min(width, height)
        left = left + (width - side) // 2
        top = top + (height - side) // 2
        width = height = side

    shape = slide.shapes.add_shape(
        shape_type,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
    )
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False  # 明示的改行がない限り1行で表示
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    # 分岐図形: 条件分岐は菱形に✕、並行は菱形に＋（DoD）。成果物・サービスは label を表示。
    if node.type == "gateway":
        p.text = "＋" if node.gateway_type == "parallel" else "✕"
    else:
        p.text = node.label  # task / start / end / artifact / service
    p.font.size = Pt(10)
    p.font.bold = False
    p.font.color.rgb = RGBColor(0, 0, 0)  # 黒文字（DoD: タスク文字の配置）
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    shape.line.color.rgb = RGBColor(0x37, 0x37, 0x37)
    # タスクの四角には影を付けない（DoD）
    shape.shadow.inherit = False
    return shape


def yaml_to_pptx(
    yaml_path: str | Path,
    output_path: str | Path,
) -> int:
    """
    YAML ファイルを読み、PPTX レイアウト仕様に従って編集可能な PPTX を生成する。
    戻り値はスライドに追加した図形の総数（タスク・分岐・矢印・レーン線・ラベル含む）。
    """
    actors, nodes = load_process_yaml(yaml_path)
    if not actors or not nodes:
        prs = Presentation()
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(6858000)
        blank = prs.slide_layouts[6]
        prs.slides.add_slide(blank)
        prs.save(str(output_path))
        return 0

    layout = compute_layout(actors, nodes)
    prs = Presentation()
    prs.slide_width = Emu(layout.slide_width)
    prs.slide_height = Emu(layout.slide_height)
    blank = prs.slide_layouts[6]

    total_shapes = 0

    for slide_idx in range(layout.num_slides):
        slide = prs.slides.add_slide(blank)
        shape_by_id = {}

        # アクター名（左）
        _draw_actor_labels(slide, layout)
        total_shapes += len(layout.actors)

        # レーン区切り（グレー点線）
        _draw_lane_separators(slide, layout)
        total_shapes += max(0, len(layout.actors) - 1)

        # このスライドに属するノード
        for node in layout.nodes:
            if node.slide_index != slide_idx:
                continue
            pos = layout.node_positions.get(node.id)
            if not pos:
                continue
            left, top, w, h = pos
            shp = _draw_node_shape(slide, node, left, top, w, h)
            shape_by_id[node.id] = shp
            total_shapes += 1

        # このスライド内のエッジのみ矢印で接続（両端が同じスライド）
        for from_id, to_id in layout.edges:
            from_node = next((n for n in layout.nodes if n.id == from_id), None)
            to_node = next((n for n in layout.nodes if n.id == to_id), None)
            if not from_node or not to_node:
                continue
            if from_node.slide_index != slide_idx or to_node.slide_index != slide_idx:
                continue
            from_shp = shape_by_id.get(from_id)
            to_shp = shape_by_id.get(to_id)
            if not from_shp or not to_shp:
                continue
            # 同一レーン内は直線、異なるレーン間は折れ曲がり（直角コネクタ）
            same_lane = from_node.actor_index == to_node.actor_index
            connector_type = (
                MSO_CONNECTOR_TYPE.STRAIGHT
                if same_lane
                else MSO_CONNECTOR_TYPE.ELBOW
            )
            conn = slide.shapes.add_connector(
                connector_type, 0, 0, 0, 0
            )
            # 接続点: 始点は下レーン同列なら下・それ以外は右、終点は同レーン同列なら上・別列なら左（DoD）
            site_from = _connection_site_from(from_node, to_node)
            site_to = _connection_site_to(from_node, to_node)
            conn.begin_connect(from_shp, site_from)
            conn.end_connect(to_shp, site_to)
            conn.line.fill.solid()
            conn.line.fill.fore_color.rgb = RGBColor(0x37, 0x37, 0x37)
            conn.line.width = Pt(1)
            conn.shadow.inherit = False  # 矢印に影を付けない（DoD）
            _add_arrow_to_connector(conn)
            total_shapes += 1

            # 分岐矢印のラベル（Yes/No 等）を矢印の近くに表示（DoD）
            edge_label = layout.edge_labels.get((from_id, to_id))
            if edge_label:
                # 座標は layout の EMU で計算し、矢印の中点付近にテキストボックスを配置
                from_pos = layout.node_positions.get(from_id)
                to_pos = layout.node_positions.get(to_id)
                if from_pos and to_pos:
                    fl, ft, fw, fh = from_pos
                    tl, tt, tw, th = to_pos
                    # 始点・終点の中心
                    fx_c = fl + fw // 2
                    fy_c = ft + fh // 2
                    tx_c = tl + tw // 2
                    ty_c = tt + th // 2
                    mx = (fx_c + tx_c) // 2
                    my = (fy_c + ty_c) // 2
                else:
                    mx = (from_shp.left + from_shp.width // 2 + to_shp.left + to_shp.width // 2) // 2
                    my = (from_shp.top + from_shp.height // 2 + to_shp.top + to_shp.height // 2) // 2
                label_w = 360000  # 約 1cm（ラベルが収まる幅）
                label_h = 120000  # 約 3mm（8pt テキスト用）
                label_left = mx - label_w // 2
                label_top = my - label_h - 60000  # 矢印の上側にオフセット
                tb = slide.shapes.add_textbox(Emu(label_left), Emu(label_top), Emu(label_w), Emu(label_h))
                tb.shadow.inherit = False
                tf = tb.text_frame
                tf.clear()
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.text = edge_label
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.alignment = PP_ALIGN.CENTER
                total_shapes += 1

        # システム接続: 点線で人⇔サービス。request=人側○・サービス側矢印下、response=下→上点線（DoD）
        for from_id, to_id, role in layout.system_edges:
            from_node = next((n for n in layout.nodes if n.id == from_id), None)
            to_node = next((n for n in layout.nodes if n.id == to_id), None)
            if not from_node or not to_node:
                continue
            if from_node.slide_index != slide_idx or to_node.slide_index != slide_idx:
                continue
            from_shp = shape_by_id.get(from_id)
            to_shp = shape_by_id.get(to_id)
            if not from_shp or not to_shp:
                continue
            same_col = from_node.col_in_slide == to_node.col_in_slide
            connector_type = MSO_CONNECTOR_TYPE.STRAIGHT if same_col else MSO_CONNECTOR_TYPE.ELBOW
            conn = slide.shapes.add_connector(connector_type, 0, 0, 0, 0)
            if role == "request":
                # 人→サービス: 人側○、サービス側矢印（下向きに接続＝サービスは上辺で受ける）
                conn.begin_connect(from_shp, CONNECTION_SITE_RIGHT)
                conn.end_connect(to_shp, CONNECTION_SITE_TOP)
                _set_connector_ends(conn, tail_oval=True, head_arrow=True)
            else:
                # response: サービス→人、下から上へ。サービス下辺→人上辺。サービス側矢印、人側○
                conn.begin_connect(from_shp, CONNECTION_SITE_BOTTOM)
                conn.end_connect(to_shp, CONNECTION_SITE_TOP)
                ln = conn.line._get_or_add_ln()
                ln.append(
                    parse_xml(
                        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
                    )
                )
                ln.append(
                    parse_xml(
                        '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="oval" w="med" len="med"/>'
                    )
                )
            _set_connector_dotted(conn)
            conn.line.fill.solid()
            conn.line.fill.fore_color.rgb = RGBColor(0x37, 0x37, 0x37)
            conn.line.width = Pt(1)
            conn.shadow.inherit = False
            total_shapes += 1

    prs.save(str(output_path))
    return total_shapes

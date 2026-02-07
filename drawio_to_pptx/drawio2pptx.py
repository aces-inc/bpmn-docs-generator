"""Drawio（mxGraphModel）XML から編集可能な図形を含む PPTX を生成する。"""

import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Drawio の 1 単位あたりの EMU（1 inch = 914400 EMU）。100 drawio 単位 ≈ 約 1 inch になるよう調整。
EMU_PER_DRAWIO_UNIT = 9144


@dataclass
class CellGeometry:
    x: float
    y: float
    width: float
    height: float


@dataclass
class ParsedCell:
    id: str
    parent: Optional[str]
    geometry: Optional[CellGeometry]
    style: str
    value: str
    vertex: bool


def _parse_style(style: str) -> dict:
    """drawio の style 文字列を key=value の辞書に分解する。"""
    result = {}
    for part in style.split(";"):
        part = part.strip()
        if "=" in part:
            k, _, v = part.partition("=")
            result[k.strip()] = v.strip()
    return result


def _parse_geometry(elem: ET.Element) -> Optional[CellGeometry]:
    g = elem.find("mxGeometry")
    if g is None:
        return None
    x = float(g.get("x") or 0)
    y = float(g.get("y") or 0)
    w = float(g.get("width") or 0)
    h = float(g.get("height") or 0)
    if w <= 0 or h <= 0:
        return None
    return CellGeometry(x=x, y=y, width=w, height=h)


def _hex_to_rgb(hex_color: str) -> Optional[RGBColor]:
    """#RRGGBB を RGBColor に。"""
    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) == 6:
        try:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        except ValueError:
            pass
    return None


def _extract_mx_graph_model_root(xml_content: str) -> ET.Element:
    """XML から mxGraphModel の root 要素を取得する。.drawio の場合は最初の diagram 内を探す。"""
    root = ET.fromstring(xml_content)
    if root.tag == "mxfile":
        diagram = root.find(".//diagram")
        if diagram is not None:
            root = diagram
    model = root.find("mxGraphModel")
    if model is not None:
        root = model
    root_elem = root.find("root")
    if root_elem is not None:
        return root_elem
    return root


def parse_cells(xml_content: str) -> list[ParsedCell]:
    """XML から vertex の mxCell を親子を考慮してパースする。"""
    root = _extract_mx_graph_model_root(xml_content)
    cells = []
    for elem in root.iter("mxCell"):
        cell_id = elem.get("id") or ""
        parent = elem.get("parent")
        style = elem.get("style") or ""
        value = elem.get("value") or ""
        geometry = _parse_geometry(elem) if elem.find("mxGeometry") is not None else None
        vertex = elem.get("vertex") == "1"
        cells.append(
            ParsedCell(
                id=cell_id,
                parent=parent,
                geometry=geometry,
                style=style,
                value=value,
                vertex=vertex,
            )
        )
    return cells


def _shape_type_from_style(style: str) -> MSO_SHAPE:
    """style から図形種別を推測する。"""
    s = _parse_style(style)
    shape = (s.get("shape") or "").lower()
    if "ellipse" in shape or "ellipse" in style.lower():
        return MSO_SHAPE.OVAL
    if "rhombus" in shape or "rhombus" in style.lower():
        return MSO_SHAPE.DIAMOND
    return MSO_SHAPE.ROUNDED_RECTANGLE


def drawio_to_pptx(
    xml_content: str,
    output_path: str | Path,
    scale: float = EMU_PER_DRAWIO_UNIT,
) -> int:
    """
    変換済み Drawio（または同 XML）から、編集可能な図形を含む PPTX を生成する。
    戻り値はスライドに追加した図形の数。
    """
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    cells = parse_cells(xml_content)
    # ルート直下の vertex のみ描画（エッジは初期スコープ外）
    root_ids = {c.id for c in cells if not c.parent or c.parent == ""}
    if not root_ids:
        root_ids = {"0", "1"}
    drawable = [
        c
        for c in cells
        if c.vertex and c.geometry and c.parent in root_ids
    ]

    for cell in drawable:
        g = cell.geometry
        shape_type = _shape_type_from_style(cell.style)
        shape = slide.shapes.add_shape(
            shape_type,
            int(g.x * scale),
            int(g.y * scale),
            int(g.width * scale),
            int(g.height * scale),
        )
        if cell.value:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = cell.value
            p.font.size = Pt(10)
        style = _parse_style(cell.style)
        fill = style.get("fillColor")
        if fill:
            rgb = _hex_to_rgb(fill)
            if rgb:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb
        stroke = style.get("strokeColor")
        if stroke:
            rgb = _hex_to_rgb(stroke)
            if rgb:
                shape.line.color.rgb = rgb

    prs.save(str(output_path))
    return len(drawable)


def drawio_file_to_pptx(drawio_path: str | Path, output_path: str | Path) -> int:
    """ .drawio ファイルを読み、PPTX に変換する。戻り値はスライドに追加した図形の数。"""
    path = Path(drawio_path)
    xml_content = path.read_text(encoding="utf-8")
    return drawio_to_pptx(xml_content, output_path)

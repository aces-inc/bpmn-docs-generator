"""yaml2pptx のテスト。"""

from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx import Presentation

from process_to_pptx import yaml2pptx


SAMPLE_YAML = """
actors:
  - A
  - B
nodes:
  - id: 1
    type: task
    actor: 0
    label: Task 1
    next: [2]
  - id: 2
    type: task
    actor: 1
    label: Task 2
    next: []
"""


def test_yaml_to_pptx_creates_file(tmp_path: Path) -> None:
    yaml_path = tmp_path / "in.yaml"
    yaml_path.write_text(SAMPLE_YAML, encoding="utf-8")
    out = tmp_path / "out.pptx"
    n = yaml2pptx.yaml_to_pptx(yaml_path, out)
    assert out.exists()
    assert out.stat().st_size > 0
    assert n > 0


def test_yaml_to_pptx_returns_shape_count(tmp_path: Path) -> None:
    yaml_path = tmp_path / "in.yaml"
    yaml_path.write_text(SAMPLE_YAML, encoding="utf-8")
    out = tmp_path / "out.pptx"
    n = yaml2pptx.yaml_to_pptx(yaml_path, out)
    # アクターラベル2 + レーン線1 + タスク2 + 矢印1 以上
    assert n >= 6


def test_yaml_to_pptx_empty_yaml(tmp_path: Path) -> None:
    yaml_path = tmp_path / "empty.yaml"
    yaml_path.write_text("actors: []\nnodes: []\n", encoding="utf-8")
    out = tmp_path / "out.pptx"
    n = yaml2pptx.yaml_to_pptx(yaml_path, out)
    assert out.exists()
    assert n == 0


def test_actor_label_vertical_center(tmp_path: Path) -> None:
    """アクター名がスイムレーンの上下中央に配置される（DoD: アクター名の位置）。"""
    yaml_path = tmp_path / "in.yaml"
    yaml_path.write_text(SAMPLE_YAML, encoding="utf-8")
    out = tmp_path / "out.pptx"
    yaml2pptx.yaml_to_pptx(yaml_path, out)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    # アクターラベルはテキストボックスで、vertical_anchor が MIDDLE であること
    textboxes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    assert len(textboxes) >= 2, "少なくとも2つのアクターラベルがある"
    for tb in textboxes[:2]:  # 最初の2つはアクター名
        assert tb.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE


SAMPLE_YAML_START_END = """
actors:
  - A
nodes:
  - id: 1
    type: start
    actor: 0
    label: 開始
    next: [2]
  - id: 2
    type: task
    actor: 0
    label: 作業
    next: [3]
  - id: 3
    type: end
    actor: 0
    label: 終了
    next: []
"""


def test_start_end_drawn_as_oval(tmp_path: Path) -> None:
    """スタート・終了ノードは PPTX 上で正円（OVAL）として描画される（DoD）。"""
    yaml_path = tmp_path / "in.yaml"
    yaml_path.write_text(SAMPLE_YAML_START_END.strip(), encoding="utf-8")
    out = tmp_path / "out.pptx"
    n = yaml2pptx.yaml_to_pptx(yaml_path, out)
    assert out.exists()
    assert n > 0
    prs = Presentation(str(out))
    assert len(prs.slides) >= 1
    slide = prs.slides[0]
    ovals = []
    for s in slide.shapes:
        try:
            if hasattr(s, "auto_shape_type") and s.auto_shape_type == MSO_SHAPE.OVAL:
                ovals.append(s)
        except (ValueError, AttributeError):
            pass  # コネクタ・テキストボックス等は auto shape ではない
    assert len(ovals) >= 2, "start と end の少なくとも2つが正円（OVAL）で描画されていること"


SAMPLE_YAML_GATEWAY_SYMBOLS = """
actors:
  - A
nodes:
  - id: 1
    type: gateway
    actor: 0
    label: 条件
    gateway_type: exclusive
    next: [2, 3]
  - id: 2
    type: gateway
    actor: 0
    label: 並行
    gateway_type: parallel
    next: [4]
  - id: 3
    type: task
    actor: 0
    label: T3
    next: [4]
  - id: 4
    type: task
    actor: 0
    label: T4
    next: []
"""


def test_gateway_drawn_with_x_or_plus(tmp_path: Path) -> None:
    """条件分岐は菱形に✕、並行分岐は菱形に＋で描画される（DoD）。"""
    yaml_path = tmp_path / "in.yaml"
    yaml_path.write_text(SAMPLE_YAML_GATEWAY_SYMBOLS.strip(), encoding="utf-8")
    out = tmp_path / "out.pptx"
    yaml2pptx.yaml_to_pptx(yaml_path, out)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    diamond_texts = []
    for s in slide.shapes:
        try:
            if hasattr(s, "auto_shape_type") and s.auto_shape_type == MSO_SHAPE.DIAMOND:
                diamond_texts.append(s.text_frame.paragraphs[0].text if s.text_frame.paragraphs else "")
        except (ValueError, AttributeError):
            pass
    assert "✕" in diamond_texts, "条件分岐（exclusive）は菱形に✕"
    assert "＋" in diamond_texts, "並行分岐（parallel）は菱形に＋"


SAMPLE_YAML_BRANCH_LABELS = """
actors:
  - A
nodes:
  - id: 1
    type: gateway
    actor: 0
    label: 成約?
    next:
      - id: 2
        label: "Yes"
      - id: 3
        label: "No"
  - id: 2
    type: task
    actor: 0
    label: 契約
    next: []
  - id: 3
    type: task
    actor: 0
    label: 見送り
    next: []
"""


def test_branch_arrow_labels_drawn(tmp_path: Path) -> None:
    """分岐矢印にラベル（Yes/No 等）が表示される（DoD）。"""
    yaml_path = tmp_path / "in.yaml"
    yaml_path.write_text(SAMPLE_YAML_BRANCH_LABELS.strip(), encoding="utf-8")
    out = tmp_path / "out.pptx"
    n = yaml2pptx.yaml_to_pptx(yaml_path, out)
    assert out.exists()
    assert n > 0
    prs = Presentation(str(out))
    slide = prs.slides[0]
    all_text = " ".join(
        s.text_frame.paragraphs[0].text
        for s in slide.shapes
        if hasattr(s, "text_frame") and s.text_frame.paragraphs
    )
    assert "Yes" in all_text, "分岐矢印ラベル Yes がスライドに含まれる"
    assert "No" in all_text, "分岐矢印ラベル No がスライドに含まれる"
